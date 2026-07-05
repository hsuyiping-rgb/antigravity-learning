import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# 1. 初始化投影片簡報並設定為 16:9 寬螢幕
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 2. 定義簡報色彩系統 (科技智慧深色流)
DARK_BLUE = RGBColor(9, 14, 23)       # 科技極深背景色：#090e17
LIGHT_BLUE = RGBColor(19, 27, 46)     # 科技深背景色：#131b2e
TEXT_GRAY = RGBColor(226, 232, 240)   # 冰川灰白字體：#e2e8f0
HIGHLIGHT_ORANGE = RGBColor(255, 94, 98) # 螢光珊瑚橘：#ff5e62
BG_LIGHT = RGBColor(9, 14, 23)        # 主背景色

AURORA_BLUE = RGBColor(0, 242, 254)   # 極光藍：#00f2fe
AURORA_CYAN = RGBColor(79, 172, 254)  # 炫光藍綠：#4facfe
CARD_DARK = RGBColor(20, 30, 48)      # 暗色卡片背景
CARD_BORDER = RGBColor(0, 242, 254)   # 卡片極光藍邊條
TEXT_MUTED = RGBColor(148, 163, 184)  # 次要字體：#94a3b8

# ----------------- 文字大小設定 (可在下方修改 Pt 值) -----------------
FONT_LEAD_TITLE = Pt(48)      # 封面/封底主標題大小
FONT_LEAD_SUBTITLE = Pt(24)   # 封面/封底副標題大小
FONT_LEAD_META = Pt(16)       # 封面/封底署名/時間大小

FONT_SLIDE_TITLE = Pt(34)     # 標準投影片單頁大標題大小
FONT_SLIDE_BODY = Pt(24)      # 標準投影片清單內文大小
FONT_SLIDE_BLOCKQUOTE = Pt(22) # 老師痛點卡片內文大小
# ------------------------------------------------------------------

# 3. 輔助函數：添加項目符號並支援 **粗體** 自動高亮
def add_bullet_point(tf, text, level=0, font_size=FONT_SLIDE_BODY):
    p = tf.add_paragraph()
    p.level = level
    p.space_after = Pt(8)
    
    # 移除 Markdown 的代碼框
    text = text.replace("`", "")
    
    # 簡單解析 **粗體**
    parts = text.split("**")
    for idx, part in enumerate(parts):
        run = p.add_run()
        run.text = part
        run.font.size = font_size
        run.font.name = "Microsoft JhengHei"
        if idx % 2 == 1:
            run.font.bold = True
            run.font.color.rgb = HIGHLIGHT_ORANGE
        else:
            run.font.color.rgb = TEXT_GRAY

# 4. 輔助函數：添加卡片式高亮區塊 (Blockquote/Pain point)
def add_blockquote(slide, text, left, top, width, height):
    # 背景圓角矩形
    bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = CARD_DARK
    bg.line.color.rgb = CARD_DARK
    
    # 左側飾條
    border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.08), height)
    border.fill.solid()
    border.fill.fore_color.rgb = CARD_BORDER
    border.line.color.rgb = CARD_BORDER
    
    # 卡片內文
    tb = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.1), width - Inches(0.3), height - Inches(0.2))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0
    p = tf.paragraphs[0]
    p.text = text.strip()
    p.font.size = FONT_SLIDE_BLOCKQUOTE
    p.font.name = "Microsoft JhengHei"
    p.font.color.rgb = TEXT_MUTED

# 5. 簡報資料定義 (15頁)
slides_data = [
    {
        "type": "lead",
        "title": "從「被動聊天」到「主動做事」",
        "image": "slides/images/slide_1.png",
        "subtitle": "教師的 AI Agent (虛擬助教) 實戰啟蒙",
        "meta": ""
    },
    {
        "type": "standard",
        "title": "傳統 AI 聊天 vs. AI Agent 助教",
        "image": "slides/images/slide_3.png",
        "bullets": [
            "傳統 AI (Chat)",
            "  - 像「點唱機」：問一句答一句，手動做所有事。",
            "  - 沒有記憶，無法幫您操作電腦或使用工具。",
            "AI Agent (虛擬助教)",
            "  - 像「實習生」：給它一個「目標」，它會自己想辦法完成。",
            "  - 能讀寫檔案、寫簡單的程式、上網找資料，還會自我除錯。",
            "  - 運作流程：查資料、提計畫、動手做、檢查結果。"
        ]
    },
    {
        "type": "standard",
        "title": "虛擬助教的準備：快速上手",
        "image": "slides/images/slide_5.png",
        "bullets": [
            "快速建置環境",
            "  - 使用 `uv` (極速 Python 套件與環境管理工具)，一鍵建立一個乾淨、不干擾系統的 Python 虛擬工作間。",
            "指定專屬工作資料夾 (Workspace)",
            "  - 告訴助教：「你只能在這個資料夾內活動與修改檔案」。",
            "  - 確保助教執行指令與檔案讀寫時的安全，防止動到電腦的其他重要檔案。"
        ]
    },
    {
        "type": "standard",
        "title": "虛擬助教的準備：設定與權限簿",
        "image": "slides/images/slide_5.png",
        "bullets": [
            "全域設定 (Global)",
            "  - 路徑：`C:\\Users\\vm\\.gemini\\config`",
            "  - 存放助教通用的密鑰、您的登入帳號與預設 AI 模型。",
            "專案設定 (Project)",
            "  - 路徑：專案資料夾下的 `.agents`",
            "  - 存放該專案專屬的工作規則，告訴助教什麼能做、什麼不能做。"
        ]
    },
    {
        "type": "standard",
        "title": "虛擬助教的準備：第一步（安裝軟體與套件）",
        "image": "slides/images/slide_5.png",
        "bullets": [
            "安裝軟體與套件",
            "  - 主程式下載：前往 [antigravity.google/download](https://antigravity.google/download) 下載安裝 Antigravity 2.0 桌面版。",
            "  - 帳戶授權：啟動後登入您的 Google 帳戶完成 AI 權限登記。"
        ]
    },
    {
        "type": "standard",
        "title": "虛擬助教的準備：第二步（建立工作目錄）",
        "image": "slides/images/slide_5.png",
        "bullets": [
            "建立工作目錄",
            "  - 建立一個全新的專案資料夾。",
            "  - 設定好版本控制與虛擬環境，讓助教隨時可以開工。"
        ]
    },
    {
        "type": "standard",
        "title": "設定面板(上)：基礎與權限守則",
        "image": "slides/images/slide_6.png",
        "bullets": [
            "定義助教的基本環境與安全限制",
            "  - General (一般)：設定命令視窗與工作資料夾。(例如：指定專案路徑，防止助教存取其他私人目錄)",
            "  - Account (帳戶)：登入 Google 帳戶啟用 AI 權限。(例如：登入您的 Google 帳號以調用 Gemini 模型)",
            "  - Permissions (權限)：最安全！ 設定執行敏感命令時需老師點擊同意。(例如：助教修改程式碼前會先彈出確認按鈕)",
            "  - Appearance (外觀)：調整介面字型與主題色彩。(例如：切換成高對比暗黑模式，字體放大至 16pt)"
        ]
    },
    {
        "type": "standard",
        "title": "設定面板(下)：進階與客製化守則",
        "image": "slides/images/slide_6.png",
        "bullets": [
            "自訂助教的智能大腦與特殊技能",
            "  - Models (模型)：挑選適合的 AI 模型作為助教的智囊。(例如：選擇高脈絡的 Gemini 1.5 Pro 來分析數小時的影片)",
            "  - Customizations (客製化)：載入特定工作流技能或規則。(例如：匯入 `classroom-video-analyzer` 技能)",
            "  - Browser (瀏覽器)：設定助教上網爬取資料時的設定。(例如：啟用無頭瀏覽器自動抓取學術論文與新聞)",
            "  - App (應用程式)：管理系統暫存與操作日誌。(例如：查看助教昨日的完整歷史運行日誌來進行偵錯)"
        ]
    },
    {
        "type": "standard",
        "title": "脈絡工具：給助教的「參考資料與任務指令」",
        "image": "slides/images/slide_7.png",
        "bullets": [
            "`Media (上傳參考資料)`：可以直接把簡報 PDF、課堂照片、影片拉給助教，多模態 AI 能直接看懂排版或 UI 設計。",
            "`@ 標記提及 (Mentions)`：精準交代資料。標記 `@ 某個檔案` 就像是「直接把考卷遞給助教看」，或指名特定的專業子助教。",
            "`Actions (快捷動作)`：輸入斜線指令（如 `/goal` 直接交代大目標，`/schedule` 交代定時自動工作）。",
            "`Browser (網頁瀏覽)`：讓助教開啟瀏覽器上網，抓取最新的時事或學術資料來設計教材。"
        ]
    },
    {
        "type": "standard",
        "title": "實戰演練：一鍵安裝 Antigravity 專屬懶人包",
        "image": "slides/images/slide_lazy_pack.png",
        "bullets": [
            "快速部署整合環境與 MCP 服務",
            "  - 懶人包下載：下載 [Antigravity 專屬懶人包](https://github.com/mathruffian-dot/antigravity-lazy-pack/tree/main) 內容。",
            "  - 懶人包安裝：將 `09-AntiGravity專屬懶人包.md` 文件提供給助教，並對助教下達「安裝這個懶人包」指令。",
            "  - 💡 懶人包自動配置服務：",
            "  - NotebookLM 整合：自動登入 Google 帳戶並完成憑證升級，確保大腦連線順暢。",
            "  - Obsidian 庫對接：自動連結並索引您的第二大腦本地知識庫。",
            "  - 專案自動化：一鍵設定 Git 版本控制與全域專案初始流程。"
        ]
    },
    {
        "type": "standard",
        "title": "什麼是 Skill？助教的「SOP 教戰手冊」",
        "image": "slides/images/slide_8.png",
        "bullets": [
            "Skill (技能) 的結構",
            "  - 它是一份寫好的標準作業程序檔案（如 `SKILL.md`），搭配輔助程式。",
            "觸發機制",
            "  - 當老師提出需求（如「我想設計教案」）時，助教若發現手冊有寫，就會自動套用 SOP。",
            "教學應用",
            "  - 老師可自訂 `lesson-plan-generator` 或期末考審題手冊，確保助教產出的講義與題目都符合學校的格式標準。"
        ]
    },
    {
        "type": "standard",
        "title": "實戰演練：SOP 技能安裝與觸發",
        "image": "slides/images/slide_skill_install.png",
        "bullets": [
            "技能安裝與工作流實作",
            "  - 技能檔放置：將自訂的 `SKILL.md` 放置於專案 `.agents/skills/<技能名稱>/` 下，系統會自動偵測載入該技能。",
            "  - 開工與收工 SOP：說「開工」時，自動讀取 Obsidian `每日筆記` 並檢查 Git；說「收工」時，自動安全掃描、Git 推送存檔。",
            "  - 專案初始化：說「初始化專案」時，自動建立檔案、Git 倉庫、GitHub 遠端儲存庫與 Obsidian 目錄。"
        ]
    },
    {
        "type": "standard",
        "title": "實戰任務一：NotebookLM 與 Antigravity 協同雙腦工作流",
        "image": "slides/images/slide_12.png",
        "blockquote": "一句話總結：NotebookLM 幫您「讀書動腦」，Antigravity 幫您「動手完成」。",
        "bullets": [
            "NotebookLM (大腦 — 負責讀書思考)",
            "  - 適合：讀入講義、課綱、PDF 檔案。",
            "  - 功能：做摘要、整理學習指引、回答老師提問、生成對答語音。",
            "Antigravity (雙手 — 負責動手建造)",
            "  - 適合：將軍師思考好的架構與大綱，化為實際行動。",
            "  - 功能：啟動 `Planning Mode` (寫企劃書)，自動寫出程式碼、生成簡報、發布網頁。"
        ]
    },
    {
        "type": "standard",
        "title": "實戰任務二：課堂影片自動剪接與字幕對齊",
        "image": "slides/images/slide_9.png",
        "blockquote": "老師痛點：下課後要把 2 小時的課堂影片剪出重點，還要手動聽寫、對齊字幕，非常耗時。",
        "bullets": [
            "虛擬助教的自動化三步驟：",
            "  - 聽寫字幕：助教聽取影片聲音，自動生成對齊時間軸的 `.srt` 字幕檔。",
            "  - 尋找精華：助教自動讀取字幕內容，找出包含關鍵教學重點的時間區段。",
            "  - 自動剪輯：助教自動寫出剪接指令，剪出重點影片片段並把字幕對齊壓好。"
        ]
    },
    {
        "type": "standard",
        "title": "實戰任務二實務：學習共同體 SLC 觀課 Skill 整合",
        "image": "slides/images/slide_slc_skill.png",
        "bullets": [
            "透過安裝 SLC 自訂 Skill 進行觀課分析",
            "  - 第一步：安裝 Skill：下載 [SLC-skill](https://github.com/hsuyiping-rgb/SLC-skill) 技能資料夾，放置於專案 `.agents/skills/` 下以載入技能。",
            "  - 第二步：提供材料並觸發：在工作區放入影片或課堂文字紀錄（如 `slc_lesson.txt`），對助教說「進行學共觀課分析」。",
            "  - 💡 實際測試實作任務：",
            "  - 操作：分析課堂記錄並自動依「學習共同體（SLC）」理論標記傾聽與協同學習歷程。",
            "  - 結果：助教自動分析學生傾聽與協同學習的起始時間、標記師生互動，並一鍵編譯輸出排版精美的 Word 觀課報告 `slc_report.docx`。"
        ]
    },
    {
        "type": "standard",
        "title": "實戰任務三：講義大綱一鍵生成精美投影片",
        "image": "slides/images/slide_10.png",
        "blockquote": "老師痛點：每次備課都要花大量時間調整簡報字型、對齊排版，無法專注在內容設計上。",
        "bullets": [
            "虛擬助教的解法：",
            "  - 寫程式生簡報：助教自動寫出 Python 程式，直接產出一份可供老師自由修改的 `.pptx` 簡報檔。",
            "  - 文字大綱秒變投影片：助教撰寫符合 `Marp` 格式的大綱，一鍵轉換成簡報網頁或 PDF。",
            "  - AI 自動畫插圖：助教自動生成適合簡報內容的教學配圖，不用再到處找無版權圖片。"
        ]
    },
    {
        "type": "standard",
        "title": "實戰任務四：隨堂測驗網頁一鍵放上雲端",
        "image": "slides/images/slide_11.png",
        "blockquote": "老師痛點：寫好一個互動測驗網頁後，不知道要怎麼放上網路給學生連線使用。",
        "bullets": [
            "虛擬助教的解法：",
            "  - 自動雲端存檔：助教自動幫您把專案做好 Git 控制，一鍵推上雲端代管平台。",
            "  - 一鍵完成網站發布：助教自動跑好 Firebase 部署流程，在幾秒鐘內把網頁推上雲端網站。",
            "教學應用：做好的測驗網頁隨即擁有網址，上課時投在螢幕上，學生拿手機掃 QR Code 就能立即連線進行隨堂小測驗！"
        ]
    },
    {
        "type": "standard",
        "title": "AI 時代的啟示：教學典範的轉移",
        "image": "slides/images/slide_14.png",
        "bullets": [
            "核心能力的轉變：",
            "  - 學生不一定要死記所有的程式碼語法，而是學習「如何定義目標」、「如何審查助教提的計畫企劃書」，以及「如何使用工具解決問題」。",
            "教師角色的轉變：",
            "  - 教師從單純的「板書抄寫與語法教學者」，升級為「引導者與決策者」。",
            "  - 老師審查學生提出的 `implementation_plan.md`，引導學生思考與修正。"
        ]
    },
    {
        "type": "lead",
        "title": "Q & A",
        "image": "slides/images/slide_15.png",
        "subtitle": "感謝各位老師的參與！\n讓我們一起用「虛擬助教」開啟更輕鬆、更有趣的教學新旅程。",
        "meta": ""
    }
]




# 6. 建構投影片
for idx, data in enumerate(slides_data):
    # 使用空白版面 (layout 6)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    if data["type"] == "lead":
        # --- 封面或封底頁 (Dark Blue Background) ---
        # 繪製全滿深藍色背景
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = DARK_BLUE
        bg.line.color.rgb = DARK_BLUE
        
        # 標題與副標題文字框
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.5), Inches(4.5))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p_title = tf.paragraphs[0]
        p_title.text = data["title"]
        p_title.font.size = FONT_LEAD_TITLE
        p_title.font.bold = True
        p_title.font.name = "Microsoft JhengHei"
        p_title.font.color.rgb = RGBColor(255, 255, 255)
        p_title.space_after = Pt(20)
        
        p_sub = tf.add_paragraph()
        p_sub.text = data["subtitle"]
        p_sub.font.size = FONT_LEAD_SUBTITLE
        p_sub.font.name = "Microsoft JhengHei"
        p_sub.font.color.rgb = AURORA_BLUE
        p_sub.space_after = Pt(30)
        
        p_meta = tf.add_paragraph()
        p_meta.text = data["meta"]
        p_meta.font.size = FONT_LEAD_META
        p_meta.font.name = "Microsoft JhengHei"
        p_meta.font.color.rgb = RGBColor(160, 174, 192)
        
        # 插入右側配圖
        if os.path.exists(data["image"]):
            slide.shapes.add_picture(data["image"], Inches(7.8), Inches(1.0), Inches(4.8), Inches(5.5))
            
    else:
        # --- 標準內容頁 (Dark Cyber Style) ---
        # 繪製全滿深藍色背景
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = DARK_BLUE
        bg.line.fill.background() # 無邊框
        
        # 投影片大標題 (依據標題長度自動適應 Pt(24-34)，以單行為原則，避免超出第二行)
        title_len = len(data["title"])
        if title_len > 25:
            slide_title_font_size = Pt(24)
        elif title_len > 20:
            slide_title_font_size = Pt(28)
        elif title_len > 15:
            slide_title_font_size = Pt(30)
        else:
            slide_title_font_size = Pt(34)

        tb_title = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.8))
        tf_title = tb_title.text_frame
        tf_title.word_wrap = True
        tf_title.margin_left = tf_title.margin_top = tf_title.margin_bottom = tf_title.margin_right = 0
        p_h1 = tf_title.paragraphs[0]
        p_h1.text = data["title"]
        p_h1.font.size = slide_title_font_size
        p_h1.font.bold = True
        p_h1.font.name = "Microsoft JhengHei"
        p_h1.font.color.rgb = AURORA_BLUE
        
        # 標題底部分割線
        divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.3), Inches(11.7), Inches(0.03))
        divider.fill.solid()
        divider.fill.fore_color.rgb = AURORA_CYAN
        divider.line.color.rgb = AURORA_CYAN
        
        # 左側內容高度自適應計算
        content_top = Inches(1.6)
        content_height = Inches(5.0)
        
        if "blockquote" in data:
            # 如果有「痛點」或「總結」，先繪製一個高亮卡片
            add_blockquote(slide, data["blockquote"], Inches(0.8), Inches(1.6), Inches(6.5), Inches(1.2))
            content_top = Inches(2.9)
            content_height = Inches(3.7)
            
        # 計算此投影片的最佳內文字型大小 (介於 Pt(20) 到 Pt(24) 之間，避免超出頁面)
        total_text_len = sum(len(b) for b in data["bullets"])
        num_bullets = len(data["bullets"])
        if total_text_len > 150 or num_bullets > 5:
            slide_body_font_size = Pt(20)
        elif total_text_len > 100 or num_bullets > 4:
            slide_body_font_size = Pt(22)
        else:
            slide_body_font_size = Pt(24)

        # 建立內文文字框
        tb_content = slide.shapes.add_textbox(Inches(0.8), content_top, Inches(6.5), content_height)
        tf_content = tb_content.text_frame
        tf_content.word_wrap = True
        tf_content.margin_left = tf_content.margin_top = tf_content.margin_bottom = tf_content.margin_right = 0
        
        # 清除第一個預設段落的邊距
        tf_content.paragraphs[0].space_after = Pt(0)
        tf_content.paragraphs[0].space_before = Pt(0)
        
        # 逐筆加入項目符號
        first_bullet = True
        for b_text in data["bullets"]:
            level = 0
            if b_text.startswith("  -") or b_text.startswith("  1.") or b_text.startswith("  2.") or b_text.startswith("  3."):
                level = 1
                b_text = b_text.replace("  -", "").replace("  1.", "1.").replace("  2.", "2.").replace("  3.", "3.").strip()
            
            if first_bullet and level == 0:
                # 複寫第一個預設段落避免產生空白行
                p = tf_content.paragraphs[0]
                p.level = 0
                p.space_after = Pt(8)
                b_text = b_text.replace("`", "")
                parts = b_text.split("**")
                for idx, part in enumerate(parts):
                    run = p.add_run()
                    run.text = part
                    run.font.size = slide_body_font_size
                    run.font.name = "Microsoft JhengHei"
                    if idx % 2 == 1:
                        run.font.bold = True
                        run.font.color.rgb = HIGHLIGHT_ORANGE
                    else:
                        run.font.color.rgb = TEXT_GRAY
                first_bullet = False
            else:
                add_bullet_point(tf_content, b_text, level, slide_body_font_size)
                
        # 插入右側配圖
        if os.path.exists(data["image"]):
            slide.shapes.add_picture(data["image"], Inches(7.8), Inches(1.6), Inches(4.8), Inches(5.0))

# 7. 儲存原生 PPTX
output_path = "g:\\我的雲端硬碟\\antigravity2\\antigravity learning\\slides\\output_v2.pptx"
os.makedirs(os.path.dirname(output_path), exist_ok=True)
try:
    prs.save(output_path)
    print(f"Native PPTX presentation successfully compiled to {output_path}!")
except PermissionError:
    fallback_path = output_path.replace("output_v2.pptx", "output_v2_new.pptx")
    prs.save(fallback_path)
    print(f"[Warning] {output_path} is currently locked (likely open in PowerPoint).")
    print(f"Successfully compiled fallback file to: {fallback_path}")
