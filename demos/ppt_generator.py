# -*- coding: utf-8 -*-
"""
AI Agent 實戰示範：簡報自動產生器 (Python-pptx)

本腳本展示了 AI Agent（如 Antigravity）如何透過 Python 程式碼，
直接讀取大綱並自動化產生一份可編輯的 Microsoft PowerPoint (.pptx) 簡報檔案。
"""

import sys

def check_pptx_library():
    """檢測 python-pptx 是否安裝"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        print("[資訊] python-pptx 已正確安裝。")
        return True
    except ImportError:
        print("[警告] 未偵測到 python-pptx 套件。")
        print("[說明] 您可以請 Agent 自動為您安裝，或者在終端機執行: pip install python-pptx")
        return False

def generate_powerpoint_presentation(output_filename="ai_agent_intro.pptx"):
    """
    建立 PowerPoint 簡報並加入自訂投影片
    """
    if not check_pptx_library():
        print("[錯誤] 無法生成 PPTX，因為 python-pptx 程式庫尚未安裝。")
        return

    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    # 1. 建立簡報物件
    prs = Presentation()
    
    # 2. 定義顏色樣式
    DARK_BLUE = RGBColor(11, 60, 93)     # #0b3c5d
    LIGHT_BLUE = RGBColor(50, 140, 193)  # #328cc1
    TEXT_GRAY = RGBColor(51, 51, 51)     # #333333
    
    # 3. 投影片 1：封面
    slide_layout = prs.slide_layouts[5] # 空白標題版面
    slide = prs.slides.add_slide(slide_layout)
    
    # 加入主標題文字框
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(9.0), Inches(2.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "從「會聊天的 AI」到「會做事的 Agent」"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    
    # 加入副標題
    p2 = tf.add_paragraph()
    p2.text = "教師的 AI Agent 實戰啟蒙課"
    p2.font.size = Pt(24)
    p2.font.color.rgb = LIGHT_BLUE
    
    # 4. 投影片 2：傳統 AI vs. AI Agent
    slide = prs.slides.add_slide(prs.slide_layouts[1]) # 標題與清單版面
    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = "傳統 AI vs. AI Agent"
    title_shape.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE
    
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    
    p = tf.paragraphs[0]
    p.text = "傳統 AI (如普通對話 ChatGPT)"
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    
    p_sub1 = tf.add_paragraph()
    p_sub1.text = "• 被動問答，無法連結外部工具"
    p_sub1.level = 1
    p_sub1.font.color.rgb = TEXT_GRAY
    
    p2 = tf.add_paragraph()
    p2.text = "AI Agent (如 Antigravity)"
    p2.font.bold = True
    p2.font.color.rgb = DARK_BLUE
    
    p_sub2 = tf.add_paragraph()
    p_sub2.text = "• 能自主規劃步驟，讀寫本機檔案並執行程式"
    p_sub2.level = 1
    p_sub2.font.color.rgb = TEXT_GRAY

    # 5. 儲存簡報
    prs.save(output_filename)
    print(f"\n[成功] 已程式化生成 PowerPoint 簡報，儲存於：{output_filename}")

if __name__ == "__main__":
    print("=== AI Agent 簡報生成 Demo ===")
    generate_powerpoint_presentation()
